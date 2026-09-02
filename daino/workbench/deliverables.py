"""Finished files: Word, Excel, PowerPoint and PDF from a workspace document.

The markdown a workspace produces is the source of truth; these are renderings
of it for people who do not read markdown. That framing decides everything here:

* **Structure survives the crossing.** Dumping markdown text into a .docx and
  calling it a Word document is the failure mode this module exists to avoid. A
  heading becomes a heading, a table becomes a table, a bullet becomes a bullet
  at the right depth, and a slide gets a title rather than a wall of asterisks.
* **The source stays authoritative.** A deliverable is regenerated, never
  hand-edited here — which is why every one records where it came from and can
  be produced again after the document changes.
* **The parsers stay optional.** ``python-docx``, ``openpyxl`` and
  ``python-pptx`` arrive with the ``daino[documents]`` extra, exactly as the
  readers do, and a missing one is reported as a named, actionable gap. PDF is
  written here directly and therefore always available: the base install can
  always produce *something* a person can open.

There is no template engine and no styling API. Clean defaults, one pass.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Literal

from daino.workbench.extraction import missing_extra_message

#: What a workspace document can be turned into.
DeliverableFormat = Literal["docx", "xlsx", "pptx", "pdf"]

SUPPORTED_FORMATS: tuple[str, ...] = ("docx", "xlsx", "pptx", "pdf")

#: Which extra provides each writer. PDF is absent on purpose — it is written by
#: this module and needs nothing installed.
_WRITERS: dict[str, str] = {
    "docx": "python-docx",
    "xlsx": "openpyxl",
    "pptx": "python-pptx",
}


class DeliverableError(ValueError):
    """Raised when a document cannot be rendered into the requested format."""


# --------------------------------------------------------------------- blocks


@dataclass
class Block:
    """One piece of a document: a heading, a paragraph, a list, or a table."""

    kind: Literal["heading", "paragraph", "bullets", "ordered", "table", "code", "quote"]
    #: Heading depth, or list nesting is carried per item instead.
    level: int = 1
    text: str = ""
    #: For lists: (depth, text) so nesting survives into Word and PowerPoint.
    items: list[tuple[int, str]] = field(default_factory=list)
    #: For tables: the header row, then the body rows.
    rows: list[list[str]] = field(default_factory=list)


_HEADING = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET = re.compile(r"^(\s*)[-*+]\s+(.*)$")
_ORDERED = re.compile(r"^(\s*)\d+[.)]\s+(.*)$")
_TABLE_DIVIDER = re.compile(r"^\s*\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)*\|?\s*$")


def parse_markdown(text: str) -> list[Block]:
    """Turn markdown into the small block model the writers render.

    Deliberately not a full CommonMark parser: what the formats below can
    express is headings, paragraphs, lists, tables and code, so parsing more
    than that would only produce structure with nowhere to go.
    """
    blocks: list[Block] = []
    lines = text.replace("\r\n", "\n").split("\n")
    index = 0
    paragraph: list[str] = []

    def flush() -> None:
        if paragraph:
            blocks.append(Block(kind="paragraph", text=" ".join(paragraph).strip()))
            paragraph.clear()

    while index < len(lines):
        line = lines[index]
        stripped = line.strip()

        if not stripped:
            flush()
            index += 1
            continue

        heading = _HEADING.match(line)
        if heading:
            flush()
            blocks.append(
                Block(kind="heading", level=len(heading.group(1)), text=_inline(heading.group(2)))
            )
            index += 1
            continue

        if stripped.startswith("```"):
            flush()
            index += 1
            body: list[str] = []
            while index < len(lines) and not lines[index].strip().startswith("```"):
                body.append(lines[index])
                index += 1
            index += 1
            blocks.append(Block(kind="code", text="\n".join(body)))
            continue

        if stripped.startswith(">"):
            flush()
            quote = []
            while index < len(lines) and lines[index].strip().startswith(">"):
                quote.append(lines[index].strip().lstrip(">").strip())
                index += 1
            blocks.append(Block(kind="quote", text=_inline(" ".join(quote))))
            continue

        if (
            stripped.startswith("|")
            and index + 1 < len(lines)
            and _TABLE_DIVIDER.match(lines[index + 1])
        ):
            flush()
            rows = [_table_row(lines[index])]
            index += 2
            while index < len(lines) and lines[index].strip().startswith("|"):
                rows.append(_table_row(lines[index]))
                index += 1
            blocks.append(Block(kind="table", rows=rows))
            continue

        bullet = _BULLET.match(line)
        ordered = _ORDERED.match(line)
        if bullet or ordered:
            flush()
            kind: Literal["bullets", "ordered"] = "bullets" if bullet else "ordered"
            items: list[tuple[int, str]] = []
            while index < len(lines):
                match = _BULLET.match(lines[index]) or _ORDERED.match(lines[index])
                if match is None:
                    break
                depth = len(match.group(1).replace("\t", "  ")) // 2
                items.append((depth, _inline(match.group(2))))
                index += 1
            blocks.append(Block(kind=kind, items=items))
            continue

        paragraph.append(_inline(stripped))
        index += 1

    flush()
    return blocks


def _table_row(line: str) -> list[str]:
    cells = line.strip().strip("|").split("|")
    return [_inline(cell.strip()) for cell in cells]


def _inline(text: str) -> str:
    """Strip the markdown a rendered document has no use for.

    Emphasis markers are removed rather than translated: carrying bold through
    four writers is a lot of machinery for a document nobody asked to be bold,
    and leaving the asterisks in is the exact "markdown dumped into Word"
    failure this module is here to prevent. Links keep their text and lose the
    URL, which is what a printed page can act on.
    """
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1 (\2)", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


def document_title(blocks: list[Block], fallback: str) -> str:
    for block in blocks:
        if block.kind == "heading" and block.level == 1:
            return block.text
    return fallback


# --------------------------------------------------------------------- render


def render(text: str, fmt: str, *, title: str = "") -> bytes:
    """Render one markdown document into ``fmt``, as bytes ready to write."""
    if fmt not in SUPPORTED_FORMATS:
        raise DeliverableError(
            f"{fmt} is not a format Daino produces. Choose one of: "
            + ", ".join(SUPPORTED_FORMATS)
        )
    blocks = parse_markdown(text)
    heading = document_title(blocks, title or "Document")
    if fmt == "docx":
        return _docx(blocks, heading)
    if fmt == "xlsx":
        return _xlsx(blocks, heading)
    if fmt == "pptx":
        return _pptx(blocks, heading)
    return _pdf(blocks, heading)


def _require(fmt: str) -> Any:
    """Import a writer, or explain which extra provides it."""
    module = {"docx": "docx", "xlsx": "openpyxl", "pptx": "pptx"}[fmt]
    try:
        return __import__(module)
    except ImportError as exc:  # pragma: no cover - depends on the install
        raise DeliverableError(missing_extra_message(_WRITERS[fmt])) from exc


def _buffer(save: Any) -> bytes:
    from io import BytesIO

    stream = BytesIO()
    save(stream)
    return stream.getvalue()


# ----------------------------------------------------------------------- docx


def _docx(blocks: list[Block], title: str) -> bytes:
    """A Word document with real headings, lists and tables."""
    _require("docx")
    from docx import Document  # noqa: PLC0415 - optional dependency

    document = Document()
    document.core_properties.title = title
    seen_title = False
    for block in blocks:
        if block.kind == "heading":
            if block.level == 1 and not seen_title:
                document.add_heading(block.text, level=0)
                seen_title = True
                continue
            document.add_heading(block.text, level=min(block.level, 9))
        elif block.kind == "paragraph":
            document.add_paragraph(block.text)
        elif block.kind == "quote":
            document.add_paragraph(block.text, style="Intense Quote")
        elif block.kind == "code":
            paragraph = document.add_paragraph()
            run = paragraph.add_run(block.text)
            run.font.name = "Consolas"
        elif block.kind in {"bullets", "ordered"}:
            style = "List Bullet" if block.kind == "bullets" else "List Number"
            for depth, item in block.items:
                # Word's built-in list styles are numbered by depth: "List
                # Bullet 2" is the second level. Beyond three it flattens,
                # which is where an outline stops being readable anyway.
                suffix = "" if depth == 0 else f" {min(depth + 1, 3)}"
                document.add_paragraph(item, style=f"{style}{suffix}")
        elif block.kind == "table" and block.rows:
            table = document.add_table(rows=len(block.rows), cols=len(block.rows[0]))
            table.style = "Light Grid Accent 1"
            for row_index, row in enumerate(block.rows):
                for cell_index, value in enumerate(row[: len(block.rows[0])]):
                    table.cell(row_index, cell_index).text = value
            document.add_paragraph()
    return _buffer(document.save)


# ----------------------------------------------------------------------- xlsx


def _xlsx(blocks: list[Block], title: str) -> bytes:
    """A workbook: one sheet per table, plus the prose as a summary sheet.

    A spreadsheet made from a document is only useful if its tables arrive as
    tables — typed, headed, and wide enough to read — so that is what this does
    with them. The narrative goes on its own sheet rather than being dropped.
    """
    _require("xlsx")
    from openpyxl import Workbook  # noqa: PLC0415 - optional dependency
    from openpyxl.styles import Alignment, Font  # noqa: PLC0415
    from openpyxl.utils import get_column_letter  # noqa: PLC0415

    workbook = Workbook()
    tables = [block for block in blocks if block.kind == "table" and block.rows]
    sheet = workbook.active
    sheet.title = "Summary"
    sheet["A1"] = title
    sheet["A1"].font = Font(bold=True, size=14)
    row_index = 3
    for block in blocks:
        if block.kind == "heading":
            sheet.cell(row=row_index, column=1, value=block.text).font = Font(bold=True)
            row_index += 1
        elif block.kind in {"paragraph", "quote"}:
            cell = sheet.cell(row=row_index, column=1, value=block.text)
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            row_index += 1
        elif block.kind in {"bullets", "ordered"}:
            for depth, item in block.items:
                sheet.cell(row=row_index, column=1, value=f"{'    ' * depth}• {item}")
                row_index += 1
    sheet.column_dimensions["A"].width = 100

    for index, block in enumerate(tables, start=1):
        header, *body = block.rows
        name = _sheet_name(_preceding_heading(blocks, block) or f"Table {index}", workbook)
        table_sheet = workbook.create_sheet(name)
        table_sheet.append(header)
        for cell in table_sheet[1]:
            cell.font = Font(bold=True)
        for row in body:
            table_sheet.append([_typed(value) for value in row])
        table_sheet.freeze_panes = "A2"
        for column in range(1, len(header) + 1):
            longest = max(
                [len(str(header[column - 1]))]
                + [len(str(row[column - 1])) for row in body if len(row) >= column]
            )
            table_sheet.column_dimensions[get_column_letter(column)].width = min(
                max(longest + 2, 10), 60
            )
    return _buffer(workbook.save)


def _typed(value: str) -> Any:
    """Numbers arrive as numbers, so a spreadsheet can add them up."""
    cleaned = value.strip().replace(",", "")
    if not cleaned:
        return ""
    for cast in (int, float):
        try:
            return cast(cleaned)
        except ValueError:
            continue
    if cleaned.endswith("%"):
        try:
            return float(cleaned[:-1]) / 100
        except ValueError:
            return value
    return value


def _preceding_heading(blocks: list[Block], target: Block) -> str:
    heading = ""
    for block in blocks:
        if block is target:
            return heading
        if block.kind == "heading":
            heading = block.text
    return heading


def _sheet_name(raw: str, workbook: Any) -> str:
    """Excel's sheet names are 31 characters and reject five punctuation marks."""
    cleaned = re.sub(r"[\\/*?:\[\]]", "-", raw).strip()[:31] or "Table"
    name, counter = cleaned, 2
    while name in workbook.sheetnames:
        suffix = f"-{counter}"
        name = f"{cleaned[: 31 - len(suffix)]}{suffix}"
        counter += 1
    return name


# ----------------------------------------------------------------------- pptx


def _pptx(blocks: list[Block], title: str) -> bytes:
    """A deck: a title slide, then one slide per section, bullets nested."""
    _require("pptx")
    from pptx import Presentation  # noqa: PLC0415 - optional dependency
    from pptx.util import Pt  # noqa: PLC0415

    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    body_layout = presentation.slide_layouts[1]

    opening = presentation.slides.add_slide(title_layout)
    opening.shapes.title.text = title
    subtitle = _first_paragraph(blocks)
    if len(opening.placeholders) > 1 and subtitle:
        opening.placeholders[1].text = subtitle

    for heading, content in _sections(blocks, title):
        slide = presentation.slides.add_slide(body_layout)
        slide.shapes.title.text = heading
        frame = slide.placeholders[1].text_frame
        frame.clear()
        first = True
        for depth, line in content[:8]:  # a slide that scrolls is not a slide
            paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
            paragraph.text = line
            paragraph.level = min(depth, 4)
            paragraph.font.size = Pt(20 if depth == 0 else 16)
            first = False
    return _buffer(presentation.save)


def _first_paragraph(blocks: list[Block]) -> str:
    for block in blocks:
        if block.kind == "paragraph":
            return block.text[:200]
    return ""


def _sections(blocks: list[Block], title: str) -> list[tuple[str, list[tuple[int, str]]]]:
    """Group the document into slides: a heading and the lines beneath it."""
    sections: list[tuple[str, list[tuple[int, str]]]] = []
    heading = ""
    content: list[tuple[int, str]] = []
    for block in blocks:
        if block.kind == "heading":
            if block.level == 1 and block.text == title and not content:
                continue
            if heading or content:
                sections.append((heading or title, content))
            heading, content = block.text, []
            continue
        if block.kind in {"bullets", "ordered"}:
            content.extend(block.items)
        elif block.kind in {"paragraph", "quote"}:
            content.append((0, block.text))
        elif block.kind == "table" and block.rows:
            content.append((0, f"{block.rows[0][0]}: {len(block.rows) - 1} rows"))
    if heading or content:
        sections.append((heading or title, content))
    return sections


# ------------------------------------------------------------------------ pdf
#
# Written here rather than through a library so the base install can always
# produce a file a person can open. The scope is deliberately narrow: one
# standard font in two weights, wrapped text, and page numbers.

_PAGE_WIDTH, _PAGE_HEIGHT = 595, 842  # A4 in points
_MARGIN = 56
_LEADING = 15


def _pdf(blocks: list[Block], title: str) -> bytes:
    lines = _pdf_lines(blocks)
    pages = _paginate(lines)
    return _pdf_bytes(pages, title)


def _pdf_lines(blocks: list[Block]) -> list[tuple[str, str, int]]:
    """Flatten the document into (style, text, indent) lines."""
    out: list[tuple[str, str, int]] = []
    for block in blocks:
        if block.kind == "heading":
            out.append(("", "", 0))
            out.append(("bold", block.text, 0))
        elif block.kind in {"paragraph", "quote"}:
            out.extend(("plain", part, 0) for part in _wrap(block.text, 92))
            out.append(("", "", 0))
        elif block.kind in {"bullets", "ordered"}:
            for position, (depth, item) in enumerate(block.items, start=1):
                marker = "•" if block.kind == "bullets" else f"{position}."
                wrapped = _wrap(f"{marker} {item}", 88 - depth * 4)
                out.extend(("plain", part, depth + 1) for part in wrapped)
            out.append(("", "", 0))
        elif block.kind == "code":
            out.extend(("plain", part, 1) for part in block.text.split("\n"))
            out.append(("", "", 0))
        elif block.kind == "table" and block.rows:
            widths = [
                max(len(str(row[index])) for row in block.rows if len(row) > index)
                for index in range(len(block.rows[0]))
            ]
            for position, row in enumerate(block.rows):
                cells = [
                    str(row[index] if index < len(row) else "").ljust(min(widths[index], 28))[:28]
                    for index in range(len(block.rows[0]))
                ]
                out.append(("bold" if position == 0 else "plain", "  ".join(cells), 0))
            out.append(("", "", 0))
    return out


def _wrap(text: str, width: int) -> list[str]:
    words, lines, current = text.split(), [], ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) > width and current:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines or [""]


def _paginate(lines: list[tuple[str, str, int]]) -> list[list[tuple[str, str, int]]]:
    per_page = (_PAGE_HEIGHT - 2 * _MARGIN) // _LEADING
    pages = [lines[index : index + per_page] for index in range(0, len(lines), per_page)]
    return pages or [[]]


def _pdf_bytes(pages: list[list[tuple[str, str, int]]], title: str) -> bytes:
    objects: list[bytes] = []

    def add(body: bytes) -> int:
        objects.append(body)
        return len(objects)

    font_plain = add(
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica /Encoding /WinAnsiEncoding >>"
    )
    font_bold = add(
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold /Encoding /WinAnsiEncoding >>"
    )
    # Reserved so pages can name their parent before it is written.
    pages_id = add(b"")
    page_ids: list[int] = []
    for number, page in enumerate(pages, start=1):
        stream = _pdf_stream(page, number, len(pages))
        content_id = add(
            b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream"
        )
        page_ids.append(
            add(
                f"<< /Type /Page /Parent {pages_id} 0 R /MediaBox [0 0 {_PAGE_WIDTH} "
                f"{_PAGE_HEIGHT}] /Resources << /Font << /F1 {font_plain} 0 R /F2 "
                f"{font_bold} 0 R >> >> /Contents {content_id} 0 R >>".encode()
            )
        )
    kids = " ".join(f"{identifier} 0 R" for identifier in page_ids)
    objects[pages_id - 1] = (
        f"<< /Type /Pages /Count {len(page_ids)} /Kids [{kids}] >>".encode()
    )
    catalog = add(f"<< /Type /Catalog /Pages {pages_id} 0 R >>".encode())
    info = add(b"<< /Title (" + _pdf_text(title).encode("latin-1", "replace") + b") >>")

    out = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for index, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{index} 0 obj\n".encode() + body + b"\nendobj\n"
    start = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root {catalog} 0 R /Info {info} 0 R >>\n"
        f"startxref\n{start}\n%%EOF\n"
    ).encode()
    return bytes(out)


def _pdf_stream(page: list[tuple[str, str, int]], number: int, total: int) -> bytes:
    parts = [b"BT\n"]
    y = _PAGE_HEIGHT - _MARGIN
    for style, text, indent in page:
        if text:
            font = b"/F2 13" if style == "bold" else b"/F1 10"
            x = _MARGIN + indent * 14
            parts.append(font + b" Tf\n")
            parts.append(f"1 0 0 1 {x} {y} Tm\n".encode())
            parts.append(b"(" + _pdf_text(text).encode("latin-1", "replace") + b") Tj\n")
        y -= _LEADING
    parts.append(b"/F1 8 Tf\n")
    parts.append(f"1 0 0 1 {_PAGE_WIDTH - _MARGIN - 40} {_MARGIN // 2} Tm\n".encode())
    parts.append(b"(" + f"{number} / {total}".encode("latin-1") + b") Tj\n")
    parts.append(b"ET")
    return b"".join(parts)


#: The base-14 PDF fonts are Latin-1, so the typographic characters a model
#: writes constantly have to be folded rather than dropped — an em dash that
#: renders as "?" is worse than one that renders as "-".
_FOLD = {
    "\u2014": "-",
    "\u2013": "-",
    "\u2018": "'",
    "\u2019": "'",
    "\u201c": '"',
    "\u201d": '"',
    "\u2026": "...",
    "\u2022": "-",
    "\u00a0": " ",
    "\u2192": "->",
}


def _pdf_text(text: str) -> str:
    for character, replacement in _FOLD.items():
        text = text.replace(character, replacement)
    return text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")


def deliverable_path(source: str, fmt: str) -> str:
    """Where a rendering of ``source`` lands, beside the document it came from."""
    path = Path(source)
    return str(path.with_suffix(f".{fmt}"))
